from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os
import io
import base64
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests
from bs4 import BeautifulSoup
import json
from datetime import datetime

app = Flask(__name__)
CORS(app)

# Configuración
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'output'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

@app.route('/api/status')
def status():
    return jsonify({
        'status': 'ok',
        'service': 'REBKA Creative Suite API',
        'version': '1.0.0'
    })

# WhatsApp Sticker Generator
@app.route('/api/stickers/generate', methods=['POST'])
def generate_sticker():
    try:
        data = request.json
        prompt = data.get('prompt', '')
        style = data.get('style', 'modern')
        
        # Crear imagen 512x512
        img = Image.new('RGBA', (512, 512), (0, 0, 0, 0))
        draw = ImageDraw.Draw(img)
        
        # Aquí iría la integración con AI para generar el sticker
        # Por ahora creamos un placeholder
        draw.rectangle([50, 50, 462, 462], fill=(255, 107, 157, 255))
        
        # Guardar
        output_path = os.path.join(OUTPUT_FOLDER, f'sticker_{datetime.now().timestamp()}.png')
        img.save(output_path, 'PNG')
        
        return jsonify({
            'success': True,
            'url': f'/api/download/{os.path.basename(output_path)}',
            'message': 'Sticker generado exitosamente'
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# Web scraping para extraer info de empresa
@app.route('/api/analyze-url', methods=['POST'])
def analyze_url():
    try:
        data = request.json
        url = data.get('url', '')
        
        # Extraer información de la URL
        response = requests.get(url, timeout=10)
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Obtener título
        title = soup.find('title')
        title_text = title.get_text() if title else ''
        
        # Obtener meta descripción
        description = soup.find('meta', attrs={'name': 'description'})
        desc_text = description.get('content', '') if description else ''
        
        # Extraer colores del CSS (simplificado)
        colors = ['#FF6B9D', '#4ECDC4', '#2C3E50', '#FFE66D']
        
        return jsonify({
            'success': True,
            'name': title_text.split('-')[0].strip(),
            'description': desc_text,
            'colors': colors,
            'suggestions': [
                {'id': '1', 'title': f'{title_text[:20]} Classic', 'description': 'Diseño clásico de marca', 'style': 'minimal'},
                {'id': '2', 'title': 'Emoji Brand', 'description': 'Emoji representativo', 'style': 'emoji'},
                {'id': '3', 'title': 'Mascota', 'description': 'Mascota animada', 'style': 'cute'},
            ]
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# Presentation Generator
@app.route('/api/presentations/generate', methods=['POST'])
def generate_presentation():
    try:
        data = request.json
        template = data.get('template', 'corporate')
        title = data.get('title', 'Presentación')
        
        # Crear PPTX
        prs = Presentation()
        
        # Slide de título
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        # Guardar
        output_path = os.path.join(OUTPUT_FOLDER, f'presentacion_{datetime.now().timestamp()}.pptx')
        prs.save(output_path)
        
        return jsonify({
            'success': True,
            'url': f'/api/download/{os.path.basename(output_path)}',
            'message': 'Presentación generada'
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# Merchandising Designer
@app.route('/api/merchandising/design', methods=['POST'])
def design_merchandise():
    try:
        data = request.json
        product = data.get('product', 'tshirt')
        
        # Crear mockup
        img = Image.new('RGBA', (800, 800), (255, 255, 255, 255))
        draw = ImageDraw.Draw(img)
        
        # Dibujar producto base
        if product == 'tshirt':
            draw.rectangle([200, 200, 600, 600], fill=(240, 240, 240, 255), outline=(0,0,0,255), width=3)
        
        output_path = os.path.join(OUTPUT_FOLDER, f'merch_{datetime.now().timestamp()}.png')
        img.save(output_path, 'PNG')
        
        return jsonify({
            'success': True,
            'url': f'/api/download/{os.path.basename(output_path)}',
            'message': 'Diseño generado'
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# Image Editor
@app.route('/api/images/edit', methods=['POST'])
def edit_image():
    try:
        if 'image' not in request.files:
            return jsonify({'success': False, 'error': 'No image provided'}), 400
        
        file = request.files['image']
        img = Image.open(file.stream).convert('RGBA')
        
        # Aplicar filtros
        filter_type = request.form.get('filter', '')
        
        if filter_type == 'grayscale':
            img = img.convert('L').convert('RGBA')
        elif filter_type == 'blur':
            img = img.filter(ImageFilter.GaussianBlur(radius=2))
        
        output_path = os.path.join(OUTPUT_FOLDER, f'edited_{datetime.now().timestamp()}.png')
        img.save(output_path, 'PNG')
        
        return jsonify({
            'success': True,
            'url': f'/api/download/{os.path.basename(output_path)}',
            'message': 'Imagen editada'
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# Download files
@app.route('/api/download/<filename>')
def download_file(filename):
    file_path = os.path.join(OUTPUT_FOLDER, filename)
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    return jsonify({'error': 'File not found'}), 404

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)